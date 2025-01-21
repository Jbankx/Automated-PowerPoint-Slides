from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Add a slide with a title and content layout
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Add a title and subtitle to the slide
title = slide.shapes.title
title.text = "Automated PowerPoint Slide"

subtitle = slide.placeholders[1]
subtitle.text = "This slide was created using python-pptx library in Python."

# Save the presentation
prs.save('automated_presentation.pptx')

print("Presentation created successfully")