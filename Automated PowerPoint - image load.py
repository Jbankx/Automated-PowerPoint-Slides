from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Add a slide with a title and content layout
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Add a title and subtitle to the slide
title = slide.shapes.title
title.text = "Automated PowerPoint Slide"

subtitle = slide.placeholders[1]
subtitle.text = "This slide was created using python-pptx library in Python."

##############################################################################
# PowerPoint 2
# Add a blank slide
slide_layout = prs.slide_layouts[5]
slide_2 = prs.slides.add_slide(slide_layout)



# Add a title and subtitle to the slide
title = slide_2.shapes.title
title.text = "Automated PowerPoint Slide with Image"

# Add an image to the slide
img_path = './Naruto.jpeg'
left = Inches(1)
top = Inches(1)
height = Inches(2)
pic = slide_2.shapes.add_picture(img_path, left, top, height=height)

######################
# Get slide dimensions (this is new optional)
slide_width = prs.slide_width
slide_height = prs.slide_height

# Get image dimensions
image_width = pic.width
image_height = pic.height

# Calculate centered position (convert to integers)
# python-pptx library requires left, top, width, and height values to be integers.
left = int((slide_width - image_width) / 2)
top = int((slide_height - image_height) / 2)

# Reposition the image
pic.left = left
pic.top = top
######################

# Save the presentation
prs.save('automated_presentation_with_image.pptx')

print("Presentation with image created successfully!")