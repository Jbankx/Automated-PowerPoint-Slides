import pandas as pd
from pptx import Presentation
from pptx.util import Inches


def save_csv_to_pptx(csv_file, pptx_file):
    # Read CSV data
    df = pd.read_csv(csv_file)

    # Create a PowerPoint presentation object
    presentation = Presentation()   

    # Add a slide to the presentation
    slide_layout = presentation.slide_layouts[5]  # Most convenient for the table load with title
    slide = presentation.slides.add_slide(slide_layout)

    # Add a title to the slide
    title = slide.shapes.title
    title.text = "Automated csv load"
 

    # Get slide dimensions (this is new optional. to help with centering the table)
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height


    # Define table dimensions 
    rows, cols = df.shape
    #left = Inches(0.5)
    #top = Inches(1.0)
    table_width = Inches(9.0)
    table_height = Inches(0.8 * rows)

    # Define table position (this is new optional. to help with centering the table)
    # python-pptx library requires left, top, width, and height values to be integers. if not error.
    left = int((slide_width - table_width) / 2)
    top = int((slide_height - table_height) / 2)
    

    # Add table to slide (important part)
    table = slide.shapes.add_table(rows + 1, cols, left, top, table_width, table_height).table

    # Set column headers
    # Starting from the 1st column (default enumerate), take each 1st cell text of the table as your columm_names
    for col_num, column_name in enumerate(df.columns):
        table.cell(0, col_num).text = column_name

    # Fill table with CSV data
    for row_num, row_data in df.iterrows():
        for col_num, cell_value in enumerate(row_data):
            table.cell(row_num + 1, col_num).text = str(cell_value)

    # Save the presentation to a file
    presentation.save(pptx_file)

# Example usage
csv_file = 'menu_items_test.csv'  # Path to your CSV file
pptx_file = 'automated_presentation_from_csv.pptx'  # Path where you want to save the PowerPoint
save_csv_to_pptx(csv_file, pptx_file)

print("Presentation from CSV data created successfully!")