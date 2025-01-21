import csv
from pptx import Presentation


# Create a presentation object
prs = Presentation()

# Read data from a CSV file
csv_file = 'menu_items_test.csv'
with open(csv_file, 'r') as file:
    reader = csv.reader(file)
    for row in reader:
        # Add a slide with a title and content layout
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Add a title and content to the slide
        title = slide.shapes.title
        title.text = row[0]

        content = slide.placeholders[1]
        content.text = row[1]    
# It looks like the code above uses Column_1 [0] as the title and column_2 as test 
# for each row (including the column names from the .csv)

# Save the presentation
prs.save('automated_presentation_from_csv_test.pptx')

print("Presentation from CSV data created successfully!")